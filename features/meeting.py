"""會議紀錄整理：解析 .docx / .pptx / .txt，整理為三段式摘要"""
import io
import logging

import db
from features.chat import simple_complete

logger = logging.getLogger("lumio.meeting")


_PROMPT = (
    "請整理這份會議紀錄／文件「{filename}」為三段式摘要，"
    "繁體中文，純文字，禁用 Markdown：\n\n"
    "🎯 結論／決議：\n"
    "（條列關鍵決議與裁示事項）\n\n"
    "📝 後續待辦：\n"
    "（條列：負責單位／人員、任務、期限）\n\n"
    "📋 簽呈摘要：\n"
    "（一段精簡公文式描述，可直接作為簽呈附件起頭）\n\n"
    "文件內容：\n{content}"
)


def extract_docx(file_bytes: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(file_bytes))
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    paras.append(cell.text.strip())
    return "\n".join(paras)


def extract_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(file_bytes))
    out = []
    for i, slide in enumerate(prs.slides, 1):
        out.append(f"[第 {i} 張]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        out.append(para.text.strip())
    return "\n".join(out)


def analyze_meeting_file(user_id: str, file_bytes: bytes, filename: str) -> str:
    """解析會議紀錄類文件並整理；同時寫入對話記憶以便後續追問。"""
    size_mb = len(file_bytes) / 1024 / 1024
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    logger.info(f"會議紀錄分析 {filename} ({size_mb:.1f}MB)")

    try:
        if ext == "docx":
            text = extract_docx(file_bytes)
        elif ext == "pptx":
            text = extract_pptx(file_bytes)
        elif ext == "txt":
            text = file_bytes.decode("utf-8", errors="replace")
        else:
            return f"⚠️ 不支援的格式：{ext}"
    except Exception as e:
        logger.exception(f"文件解析失敗 {filename}: {e}")
        return f"⚠️ 文件解析失敗：{e}"

    if not text.strip():
        return "⚠️ 文件內容為空或無法解析"

    reply = simple_complete(_PROMPT.format(filename=filename, content=text[:15000]), max_tokens=1500)
    db.save_message(user_id, "user", f"[📄 上傳會議紀錄：{filename}（{size_mb:.1f}MB）]")
    db.save_message(user_id, "assistant", reply)
    return reply
